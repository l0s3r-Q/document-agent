"""parser.py —— 解析 .pptx：各页形状文本/结构。"""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


def parse(path: str) -> dict:
    """解析 pptx，返回各页结构（标题/正文要点/表格）。"""
    prs = Presentation(path)
    result = {"file": path, "slide_size": {
        "width_inches": round(prs.slide_width / 914400, 2),
        "height_inches": round(prs.slide_height / 914400, 2),
    }, "slides": []}

    for i, slide in enumerate(prs.slides, 1):
        entry = {"index": i, "shapes": []}
        # 占位符标题（若有；0.6.x 无标题占位符时 shapes.title 抛 ValueError）
        title = None
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title = title_shape.text_frame.text.strip()
        except ValueError:
            title = None
        if title:
            entry["title"] = title

        for shape in slide.shapes:
            stype = str(shape.shape_type) if shape.shape_type else "unknown"
            item = {"shape_type": stype.split(" (")[0] if " (" in stype else stype}
            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                if texts:
                    item["texts"] = texts
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                item["table"] = rows
            if getattr(shape, "shape_type", None) and "PICTURE" in str(shape.shape_type):
                item["picture"] = True
            entry["shapes"].append(item)
        result["slides"].append(entry)
    return result
