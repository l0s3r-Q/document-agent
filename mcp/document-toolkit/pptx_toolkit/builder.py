"""builder.py —— 按 PptxSpec 生成 .pptx（可编辑原生形状，python-pptx）。"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .themes import CN_FONT, DEFAULT_THEME, EN_FONT, THEMES

_SLIDE_SIZE = {
    "16:9": (Inches(13.333), Inches(7.5)),
    "4:3": (Inches(10), Inches(7.5)),
}

# 版式布局（按画布宽度适配）
_LAYOUT = {
    "16:9": {"bar_w": 13.333, "content_left": 1.2, "content_w": 10.9,
             "col1_left": 1.0, "col2_left": 6.9, "col_w": 5.4, "cover_w": 10.9},
    "4:3":  {"bar_w": 10.0, "content_left": 0.9, "content_w": 8.2,
             "col1_left": 0.7, "col2_left": 5.1, "col_w": 4.2, "cover_w": 8.2},
}


def _hex(color: str) -> RGBColor:
    return RGBColor.from_string(color.lstrip("#"))


def _set_font(run, name: str, size: float, bold: bool, color: RGBColor):
    """设置 run 字体（latin + eastAsia 都写，中文才能正确显示）。"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = EN_FONT
    rpr = run._r.get_or_add_rPr()
    # 删除旧 ea 后追加
    for tag in ("a:ea", "a:cs"):
        for el in rpr.findall(qn(tag)):
            rpr.remove(el)
    ea = rpr.makeelement(qn("a:ea"), {"typeface": name})
    rpr.append(ea)


def _add_text_box(slide, left: float, top: float, width: float, height: float,
                  text: str, size: float, bold: bool, color: RGBColor,
                  align=PP_ALIGN.LEFT) -> None:
    """添加文本框（可编辑原生形状）。"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, CN_FONT, size, bold, color)
    return box


def _add_bullets(slide, left: float, top: float, width: float, height: float,
                 items: list, size: float, color: RGBColor, line_spacing: float = 1.3) -> None:
    """添加要点列表（每项一个段落，支持"子要点"缩进：以 > 前缀）。"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if not isinstance(item, str):
            item = str(item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(size * 0.6)
        p.line_spacing = line_spacing
        item = str(item or "")
        sub = item.startswith(">")
        text = item.lstrip("> ").strip()
        if sub:
            p.level = 1
        run = p.add_run()
        run.text = ("• " if not sub else "– ") + text
        _set_font(run, CN_FONT, size if not sub else size - 2, False, color)


def _add_title_bar(slide, title: str, theme: dict, bar_w: float = 13.333):
    """页面上方标题条（主色块 + 标题文字，超长自动缩字号）。"""
    # 主色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(bar_w), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex(theme["primary"])
    bar.line.fill.background()
    # 标题（白色固定，深色标题条上恒可读；超 24 字符降为 24pt，超 40 字符降为 20pt）
    size = 28 if len(title) <= 24 else (24 if len(title) <= 40 else 20)
    _add_text_box(slide, 0.6, 0.15, 12, 0.7, title, size, True, _hex("FFFFFF"))


def build(spec: dict, output_path: str) -> dict:
    """按 PptxSpec 生成 pptx。

    spec 契约:
      {title, subtitle?, theme?("corporate"|"academic"|"launch"|"minimal"),
       size?("16:9"|"4:3"), author?,
       slides: [
         {type: "cover", title, subtitle?},
         {type: "agenda", items: [..]},
         {type: "section", title},
         {type: "content", title, bullets: [..]},          # "> 文本"=子要点
         {type: "two_column", title, left:{title,bullets}, right:{title,bullets}},
         {type: "table", title, rows:[[..]], header_row?(bool)},
         {type: "image", title, path, caption?},
         {type: "closing", title, subtitle?},
       ]}
    """
    if not output_path.lower().endswith(".pptx"):
        return {"ok": False, "error": "输出路径必须以 .pptx 结尾"}

    theme_name = spec.get("theme") or DEFAULT_THEME
    if theme_name not in THEMES:
        theme_name = DEFAULT_THEME  # 非法主题回退并归一化返回值
    theme = THEMES[theme_name]
    size = _SLIDE_SIZE.get(spec.get("size") or "16:9", _SLIDE_SIZE["16:9"])

    prs = Presentation()
    prs.slide_width, prs.slide_height = size
    slides = spec.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return {"ok": False, "error": "slides 必须是非空数组"}
    lay = _LAYOUT.get(spec.get("size") or "16:9", _LAYOUT["16:9"])
    warnings = []

    blank = prs.slide_layouts[6]

    for s in slides:
        stype = s.get("type", "content")
        slide = prs.slides.add_slide(blank)
        # 深色背景主题：设置页面背景，保证浅色文字可读
        if theme["background"].lower() != "ffffff":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _hex(theme["background"])

        if stype == "cover":
            cover_title = s.get("title") or spec.get("title") or ""
            cover_sub = s.get("subtitle") or spec.get("subtitle") or ""
            # 标题默认黑色；slide 级 title_color 可覆盖（如 "#1F4E79"）
            tcolor = _hex(s.get("title_color") or theme.get("title", "000000"))
            _add_text_box(slide, lay["content_left"], 2.2, lay["cover_w"], 1.6, cover_title,
                          theme["title_size"] + 4, True, tcolor, PP_ALIGN.CENTER)
            if cover_sub:
                _add_text_box(slide, lay["content_left"], 3.8, lay["cover_w"], 0.8, cover_sub,
                              20, False, _hex(theme["text"]), PP_ALIGN.CENTER)
            if spec.get("author"):
                _add_text_box(slide, lay["content_left"], 6.2, lay["cover_w"], 0.6, spec["author"],
                              14, False, _hex(theme["secondary"]), PP_ALIGN.CENTER)

        elif stype == "agenda":
            _add_title_bar(slide, "目录", theme, lay["bar_w"])
            items = []
            for i, it in enumerate(s.get("items", []), 1):
                items.append(f"{i}. {it}")
            _add_bullets(slide, lay["content_left"], 1.8, lay["content_w"], 4.8, items, 22, _hex(theme["text"]))

        elif stype == "section":
            tcolor = _hex(s.get("title_color") or theme.get("title", "000000"))
            # 左侧色带装饰
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                          Inches(0.25), Inches(7.5))
            band.fill.solid()
            band.fill.fore_color.rgb = _hex(theme["primary"])
            band.line.fill.background()
            # 编号圆点（可选 "index": "01"）
            idx = s.get("index", "")
            if idx:
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lay["content_left"]), Inches(2.1),
                                             Inches(1.1), Inches(1.1))
                dot.fill.solid()
                dot.fill.fore_color.rgb = _hex(theme["accent"])
                dot.line.fill.background()
                dot.shadow.inherit = False
                _add_text_box(slide, lay["content_left"], 2.3, 1.1, 0.8, idx,
                              30, True, _hex(theme["text"]), PP_ALIGN.CENTER)
            _add_text_box(slide, lay["content_left"], 3.4, lay["cover_w"], 1.4, s.get("title", ""),
                          theme["title_size"], True, tcolor, PP_ALIGN.CENTER)
            # 强调线
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.6), 5.0, Inches(4.0), 0.08)
            line.fill.solid()
            line.fill.fore_color.rgb = _hex(theme["accent"])
            line.line.fill.background()

        elif stype == "stats":
            # 大数字数据卡片：stats: [{value, label, sub?, color?}]
            _add_title_bar(slide, s.get("title", ""), theme, lay["bar_w"])
            items = s.get("stats", [])
            n = max(len(items), 1)
            gap, margin = 0.4, lay["content_left"]
            avail = lay["content_w"]
            card_w = (avail - gap * (n - 1)) / n if n > 1 else avail
            colors = ["EBF8FF", "F0FFF4", "FFF5F5", "FAF5FF"]
            for i, it in enumerate(items):
                x = margin + i * (card_w + gap)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(x), Inches(1.8), Inches(card_w), Inches(3.6))
                card.adjustments[0] = 0.06
                card.fill.solid()
                card.fill.fore_color.rgb = _hex(it.get("bg") or colors[i % len(colors)])
                card.line.color.rgb = _hex(theme["secondary"])
                card.line.width = Pt(0.75)
                card.shadow.inherit = False
                val = str(it.get("value", ""))
                vcolor = _hex(it.get("color") or theme["primary"])
                _add_text_box(slide, x, 2.3, card_w, 1.3, val,
                              max(36, 56 - max(len(val) - 1, 0) * 4), True, vcolor, PP_ALIGN.CENTER)
                _add_text_box(slide, x, 3.9, card_w, 0.6, it.get("label", ""),
                              18, True, _hex(theme["text"]), PP_ALIGN.CENTER)
                if it.get("sub"):
                    _add_text_box(slide, x, 4.6, card_w, 0.6, it["sub"],
                                  13, False, _hex(theme["secondary"]), PP_ALIGN.CENTER)

        elif stype == "content":
            _add_title_bar(slide, s.get("title", ""), theme, lay["bar_w"])
            if s.get("cards"):
                # 卡片化布局：cards: [{title, bullets: [], color?}]
                cs = s["cards"]
                n = max(len(cs), 1)
                gap, margin = 0.4, lay["content_left"]
                avail = lay["content_w"]
                card_w = (avail - gap * (n - 1)) / n if n > 1 else avail
                colors = ["2E75B6", "2F855A", "C53030", "6B46C1", "B7791F"]
                for i, c in enumerate(cs):
                    x = margin + i * (card_w + gap)
                    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                  Inches(x), Inches(1.6), Inches(card_w), Inches(5.0))
                    card.adjustments[0] = 0.04
                    card.fill.solid()
                    card.fill.fore_color.rgb = _hex("F7FAFC")
                    card.line.color.rgb = _hex("E2E8F0")
                    card.line.width = Pt(0.75)
                    card.shadow.inherit = False
                    # 顶部色条
                    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                 Inches(x), Inches(1.6), Inches(card_w), Inches(0.12))
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = _hex(c.get("color") or colors[i % len(colors)])
                    bar.line.fill.background()
                    _add_text_box(slide, x + 0.3, 2.0, card_w - 0.6, 0.6,
                                  c.get("title", ""), 20, True, _hex(theme["text"]))
                    _add_bullets(slide, x + 0.3, 2.8, card_w - 0.6, 3.6,
                                 c.get("bullets", []), theme["body_size"] - 3, _hex(theme["text"]))
            else:
                _add_bullets(slide, lay["content_left"], 1.6, lay["content_w"], 5.2, s.get("bullets", []),
                             theme["body_size"], _hex(theme["text"]))

        elif stype == "two_column":
            _add_title_bar(slide, s.get("title", ""), theme, lay["bar_w"])
            left, right = s.get("left", {}), s.get("right", {})
            _add_text_box(slide, lay["col1_left"], 1.6, lay["col_w"], 0.6, left.get("title", ""),
                          20, True, _hex(theme["primary"]))
            _add_bullets(slide, lay["col1_left"], 2.3, lay["col_w"], 4.4, left.get("bullets", []),
                         theme["body_size"] - 2, _hex(theme["text"]))
            _add_text_box(slide, lay["col2_left"], 1.6, lay["col_w"], 0.6, right.get("title", ""),
                          20, True, _hex(theme["primary"]))
            _add_bullets(slide, lay["col2_left"], 2.3, lay["col_w"], 4.4, right.get("bullets", []),
                         theme["body_size"] - 2, _hex(theme["text"]))

        elif stype == "table":
            _add_title_bar(slide, s.get("title", ""), theme, lay["bar_w"])
            rows = s.get("rows") or []
            if rows:
                ncols = max(len(r) for r in rows) or 1
                shape = slide.shapes.add_table(len(rows), ncols, Inches(lay["content_left"]), Inches(1.8),
                                               Inches(lay["content_w"]), Inches(4.2))
                tbl = shape.table
                for ri, row in enumerate(rows):
                    for ci in range(ncols):
                        cell = tbl.cell(ri, ci)
                        v = row[ci] if ci < len(row) else ""
                        cell.text = str(v) if v is not None else ""
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                _set_font(run, CN_FONT, 14,
                                          ri == 0 and s.get("header_row", True),
                                          _hex(theme["background"] if ri == 0 else theme["text"]))
                        if ri == 0 and s.get("header_row", True):
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _hex(theme["primary"])

        elif stype == "image":
            _add_title_bar(slide, s.get("title", ""), theme, lay["bar_w"])
            img_path = s.get("path")
            if img_path and os.path.exists(img_path):
                try:
                    slide.shapes.add_picture(img_path, Inches(lay["content_left"]), Inches(1.8),
                                             width=Inches(lay["content_w"] * 0.8))
                except Exception as e:  # noqa: BLE001  图片格式问题保留标题页，记录警告
                    warnings.append(f"第{s['index'] if 'index' in s else '?'}页图片添加失败: {type(e).__name__}")
            elif img_path:
                warnings.append(f"图片路径不存在: {img_path}")
            if s.get("caption"):
                _add_text_box(slide, lay["content_left"], 6.4, lay["content_w"], 0.5, s["caption"],
                              14, False, _hex(theme["secondary"]))

        elif stype == "chart":
            # 原生图表：chart_type: column|bar|line|pie; categories; series[{name, values}]
            _add_title_bar(slide, s.get("title", ""), theme, lay["bar_w"])
            chart_type = {
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE_MARKERS,
                "pie": XL_CHART_TYPE.PIE,
            }.get(s.get("chart_type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
            cats = s.get("categories", [])
            series = s.get("series", [])
            if cats and series:
                cd = CategoryChartData()
                cd.categories = cats
                for sr in series:
                    cd.add_series(sr.get("name", ""), tuple(sr.get("values", [])))
                gf = slide.shapes.add_chart(chart_type,
                                            Inches(lay["content_left"]), Inches(1.6),
                                            Inches(lay["content_w"]), Inches(4.8), cd)
                chart = gf.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

        elif stype == "closing":
            tcolor = _hex(s.get("title_color") or theme.get("title", "000000"))
            _add_text_box(slide, lay["content_left"], 2.6, lay["cover_w"], 1.4, s.get("title", "谢谢"),
                          theme["title_size"] + 4, True, tcolor, PP_ALIGN.CENTER)
            if s.get("subtitle"):
                _add_text_box(slide, lay["content_left"], 4.1, lay["cover_w"], 0.8, s["subtitle"],
                              18, False, _hex(theme["text"]), PP_ALIGN.CENTER)

        else:
            _add_title_bar(slide, s.get("title", ""), theme, lay["bar_w"])
            _add_bullets(slide, lay["content_left"], 1.6, lay["content_w"], 5.2, s.get("bullets", []),
                         theme["body_size"], _hex(theme["text"]))

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    result = {"ok": True, "path": output_path, "slides": len(slides), "theme": theme_name}
    if warnings:
        result["warnings"] = warnings
    return result
