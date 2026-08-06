"""quality.py —— 交付质量检查引擎：AIGC 痕迹/占位符/排版/结构体检。

目标：交付物不用修改、没有 AIGC 痕迹、没有细节错误。
检查项分 error（必须修复）与 warning（建议修复）。
"""

from __future__ import annotations

import os
import re

from .parser import parse as parse_docx

# ── AIGC 痕迹词（中文正式文档中的 AI 腔/空洞套话）──
AIGC_PATTERNS = [
    (r"总而言之|综上所述|总的来说", "AI 腔总结套话"),
    (r"值得注意的是|需要注意的是|需要指出的是", "AI 腔提醒套话"),
    (r"毋庸置疑|众所周知|不言而喻", "空洞论断"),
    (r"作为(一个)?人工智能|作为AI|作为大模型", "AI 身份自述"),
    (r"首先[，,].*其次[，,].*最后[，,]", "过度模板化排比"),
    (r"不仅[^。]{0,12}而且", "AI 腔递进句式"),
    (r"在当今(社会|时代|背景下)", "AI 腔开篇套话"),
    (r"随着(社会|时代|科技)的(发展|进步)", "AI 腔背景套话"),
]

# ── 占位符/未完成标记 ──
PLACEHOLDER_PATTERNS = [
    (r"\{\s*[a-zA-Z_\u4e00-\u9fa5]+\s*\}", "变量占位符未替换"),
    (r"【?待(补充|完善|填写)|TODO|FIXME|XXX{2,}", "未完成标记"),
    (r"此处(填写|插入|添加|补充)", "内容占位提示"),
]

# ── 排版细节 ──
EMOJI_RE = re.compile(
    "[\U0001F000-\U0001FAFF\u2600-\u27BF\uFE0F\u2190-\u21FF\u2B00-\u2BFF]"
)


def _check_text(text: str, issues: list, where: str, check_sentence_end: bool = True) -> None:
    """文本级检查：AIGC 痕迹/占位符/emoji/重复标点。"""
    if not text or not text.strip():
        return
    for pat, label in AIGC_PATTERNS:
        for m in re.finditer(pat, text):
            issues.append({
                "level": "warning", "type": "aigc_trace",
                "location": where, "detail": f"{label}：『{m.group(0)[:30]}』",
            })
    for pat, label in PLACEHOLDER_PATTERNS:
        for m in re.finditer(pat, text):
            issues.append({
                "level": "error", "type": "placeholder",
                "location": where, "detail": f"{label}：『{m.group(0)[:30]}』",
            })
    if EMOJI_RE.search(text):
        issues.append({
            "level": "warning", "type": "emoji",
            "location": where, "detail": "正文含 emoji（正式文档建议去除）",
        })
    if re.search(r"[!?！？]{3,}", text):
        issues.append({
            "level": "warning", "type": "punctuation",
            "location": where, "detail": "连续感叹/问号（情绪化表达）",
        })
    # 句末标点仅检查正文段落（标题/装饰行不要求）
    if check_sentence_end and text[-1] not in "。；：！？；;!?…—":
        # 豁免：日期行/纯数字结尾的元信息行（如 "2026年7月31日"）
        if not re.search(r"(20\d{2}年\d{1,2}月\d{1,2}日|\d{1,2}:\d{2}|[0-9]{2,4})$", text):
            issues.append({
                "level": "warning", "type": "sentence_end",
                "location": where, "detail": "段落未以句号等标点结尾",
            })


def check_docx(path: str) -> dict:
    """docx 质量体检：内容痕迹 + 排版 + 结构 + 重复段落 + 数据一致性。"""
    data = parse_docx(path)
    issues = []
    texts_checked = 0
    structure = data.get("structure", [])
    # ── 重复段落检测（AI 常见：整段内容高度重复/复读）──
    prev_texts = []
    for i, item in enumerate(structure):
        stype = item.get("type", "paragraph")
        text = item.get("text", "")
        if stype in ("paragraph", "title") and text and len(text.strip()) >= 12:
            norm = re.sub(r"\s+", "", text)
            # 与最近 8 段比较，相似度 > 0.9 视为重复
            for pi, prev in enumerate(prev_texts):
                if not prev or abs(len(norm) - len(prev)) > max(4, len(norm) // 5):
                    continue
                # 简单相似度：公共前缀比例
                common = 0
                for a, b in zip(norm, prev):
                    if a == b:
                        common += 1
                    else:
                        break
                if common > 0 and common / max(len(norm), len(prev)) > 0.9:
                    issues.append({
                        "level": "warning", "type": "dup_paragraph",
                        "location": f"第{i + 1}段",
                        "detail": f"与前文段落重复度高（相似前缀 {common} 字）",
                    })
                    break
            prev_texts.append(norm)
            if len(prev_texts) > 8:
                prev_texts.pop(0)

    for i, item in enumerate(structure):
        stype = item.get("type", "paragraph")
        text = item.get("text", "")
        is_heading = stype in ("heading1", "heading2", "heading3")
        is_decor = stype in ("title", "separator")
        if stype in ("paragraph", "title", "heading1", "heading2", "heading3",
                     "separator", "list", "table"):
            if text and text.strip():
                texts_checked += 1
                # 标题/装饰行不要求句末标点
                _check_text(text, issues, f"第{i + 1}段",
                            check_sentence_end=not (is_heading or is_decor))
        if stype == "table":
            rows = item.get("rows", [])
            if rows:
                ncols = max(len(r) for r in rows)
                for ri, r in enumerate(rows):
                    if len(r) < ncols:
                        issues.append({
                            "level": "error", "type": "table_ragged",
                            "location": f"第{i + 1}段表格第{ri + 1}行",
                            "detail": f"列数不足（{len(r)}/{ncols}）",
                        })
                if any(not str(c).strip() for c in rows[0]):
                    issues.append({
                        "level": "warning", "type": "table_header_empty",
                        "location": f"第{i + 1}段表格", "detail": "表头含空单元格",
                    })
        if stype == "paragraph" and text.strip():
            if not item.get("east_asia") and not item.get("font_name"):
                issues.append({
                    "level": "warning", "type": "font_missing",
                    "location": f"第{i + 1}段", "detail": "未检测到字体设置",
                })

    # 结构检查：标题层级跳级
    prev_level = 0
    for i, item in enumerate(structure):
        t = item.get("type", "")
        if t in ("heading1", "heading2", "heading3"):
            level = int(t[-1])
            if prev_level and level > prev_level + 1:
                issues.append({
                    "level": "warning", "type": "heading_skip",
                    "location": f"第{i + 1}段", "detail": f"标题层级跳级（{prev_level}→{level}）",
                })
            prev_level = level

    return {
        "ok": True, "path": path, "texts_checked": texts_checked,
        "error_count": sum(1 for x in issues if x["level"] == "error"),
        "warning_count": sum(1 for x in issues if x["level"] == "warning"),
        "issues": issues,
        "pass": all(x["level"] != "error" for x in issues),
    }


def check_pptx(path: str) -> dict:
    """pptx 质量体检：空页/文本溢出估算。"""
    from pptx import Presentation
    prs = Presentation(path)
    issues = []
    for i, slide in enumerate(prs.slides, 1):
        has_text = False
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                has_text = True
                text = sh.text_frame.text
                width_in = sh.width / 914400 if sh.width else 0
                if width_in > 0:
                    est_chars_per_line = max(int(width_in * 2.54 * 2 * 72 / 18), 10)
                    if len(text) > est_chars_per_line * 6:
                        issues.append({
                            "level": "warning", "type": "text_overflow",
                            "location": f"第{i}页", "detail": f"文本可能溢出（{len(text)} 字符）",
                        })
        if not has_text:
            issues.append({
                "level": "warning", "type": "empty_slide",
                "location": f"第{i}页", "detail": "空白页（无文本）",
            })
    return {
        "ok": True, "path": path, "slides": len(prs.slides),
        "error_count": sum(1 for x in issues if x["level"] == "error"),
        "warning_count": sum(1 for x in issues if x["level"] == "warning"),
        "issues": issues,
        "pass": all(x["level"] != "error" for x in issues),
    }


def check_excel(path: str) -> dict:
    """xlsx 质量体检：表头重复/空列/空行/数据重复/空行比例。"""
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    issues = []
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows or all(v is None for v in rows[0]):
            issues.append({
                "level": "warning", "type": "empty_header",
                "location": f"sheet「{ws.title}」", "detail": "表头为空",
            })
            continue
        header = [str(v).strip() if v is not None else "" for v in rows[0]]
        seen = {}
        for hi, h in enumerate(header):
            if h:
                if h in seen:
                    issues.append({
                        "level": "warning", "type": "dup_header",
                        "location": f"sheet「{ws.title}」", "detail": f"重复表头『{h}』",
                    })
                seen[h] = True
            else:
                issues.append({
                    "level": "warning", "type": "empty_header_cell",
                    "location": f"sheet「{ws.title}」第{hi + 1}列", "detail": "表头单元格为空",
                })
        # ── 数据质量：空行比例 + 纯重复行 ──
        data_rows = rows[1:]
        if data_rows:
            empty_rows = [r for r in data_rows if all(v is None or str(v).strip() == "" for v in r)]
            empty_ratio = len(empty_rows) / len(data_rows)
            if len(data_rows) >= 5 and empty_ratio > 0.5:
                issues.append({
                    "level": "warning", "type": "many_empty_rows",
                    "location": f"sheet「{ws.title}」",
                    "detail": f"空行占比过高（{len(empty_rows)}/{len(data_rows)}，{empty_ratio:.0%}）",
                })
            # 纯重复行检测（整行内容完全一致）
            seen_rows = set()
            dup_count = 0
            for r in data_rows:
                if all(v is None or str(v).strip() == "" for v in r):
                    continue
                key = tuple(str(v).strip() if v is not None else "" for v in r)
                if key in seen_rows:
                    dup_count += 1
                else:
                    seen_rows.add(key)
            if dup_count >= 3:
                issues.append({
                    "level": "warning", "type": "dup_rows",
                    "location": f"sheet「{ws.title}」",
                    "detail": f"检测到 {dup_count} 行完全重复的数据行",
                })
    return {
        "ok": True, "path": path, "sheets": len(wb.sheetnames),
        "error_count": sum(1 for x in issues if x["level"] == "error"),
        "warning_count": sum(1 for x in issues if x["level"] == "warning"),
        "issues": issues,
        "pass": all(x["level"] != "error" for x in issues),
    }


def quality_check(path: str) -> dict:
    """按扩展名分派质量体检。"""
    if not os.path.exists(path):
        return {"ok": False, "error": f"文件不存在: {path}"}
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        return check_docx(path)
    if ext == ".pptx":
        return check_pptx(path)
    if ext == ".xlsx":
        return check_excel(path)
    return {"ok": False, "error": f"暂不支持的类型: {ext}（支持 docx/pptx/xlsx）"}
